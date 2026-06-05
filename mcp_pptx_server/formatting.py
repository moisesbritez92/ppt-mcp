from typing import Any, Dict, Union, Optional
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from mcp_pptx_server.presentation_editor import find_shape

def set_shape_font(prs: Presentation, slide_index: int, shape_name_or_index: Union[str, int], font_name: str = "", font_size: Optional[float] = None, bold: Optional[bool] = None, italic: Optional[bool] = None, color_hex: str = "") -> Dict[str, Any]:
    """Applies font formatting to a shape (its text frame paragraphs and runs)."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"Slide index {slide_index} is out of bounds (0 to {len(prs.slides) - 1}).")

    slide = prs.slides[slide_index]
    shape = find_shape(slide, shape_name_or_index)

    if not shape.has_text_frame:
        raise ValueError(f"Shape '{shape.name}' does not support text formatting (no text frame).")

    # Parse color
    rgb_color = None
    if color_hex:
        clean_hex = color_hex.lstrip('#')
        if len(clean_hex) == 6:
            try:
                r = int(clean_hex[0:2], 16)
                g = int(clean_hex[2:4], 16)
                b = int(clean_hex[4:6], 16)
                rgb_color = RGBColor(r, g, b)
            except ValueError:
                raise ValueError(f"Invalid hex color value: {color_hex}")
        else:
            raise ValueError("color_hex must be in RRGGBB or #RRGGBB format.")

    # Apply styling
    for paragraph in shape.text_frame.paragraphs:
        if not paragraph.runs:
            # Stylize paragraph defaults if empty
            font = paragraph.font
            if font_name:
                font.name = font_name
            if font_size is not None:
                font.size = Pt(font_size)
            if bold is not None:
                font.bold = bold
            if italic is not None:
                font.italic = italic
            if rgb_color is not None:
                font.color.rgb = rgb_color
        else:
            # Stylize individual runs
            for run in paragraph.runs:
                font = run.font
                if font_name:
                    font.name = font_name
                if font_size is not None:
                    font.size = Pt(font_size)
                if bold is not None:
                    font.bold = bold
                if italic is not None:
                    font.italic = italic
                if rgb_color is not None:
                    font.color.rgb = rgb_color

    return {
        "slide_index": slide_index,
        "shape_name": shape.name,
        "shape_id": shape.shape_id,
        "font_name": font_name or "Unchanged",
        "font_size": font_size,
        "bold": bold,
        "italic": italic,
        "color_hex": color_hex or "Unchanged"
    }

def set_slide_background(prs: Presentation, slide_index: int, color_hex: str) -> Dict[str, Any]:
    """Applies a solid background color to a slide."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"Slide index {slide_index} is out of bounds (0 to {len(prs.slides) - 1}).")

    slide = prs.slides[slide_index]

    # Parse color
    clean_hex = color_hex.lstrip('#')
    if len(clean_hex) == 6:
        try:
            r = int(clean_hex[0:2], 16)
            g = int(clean_hex[2:4], 16)
            b = int(clean_hex[4:6], 16)
            rgb_color = RGBColor(r, g, b)
        except ValueError:
            raise ValueError(f"Invalid hex color value: {color_hex}")
    else:
        raise ValueError("color_hex must be in RRGGBB or #RRGGBB format.")

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color

    return {
        "slide_index": slide_index,
        "background_color_hex": color_hex
    }
