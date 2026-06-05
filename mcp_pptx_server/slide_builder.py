import os
from typing import Any, Dict, List, Optional
from pptx import Presentation
from pptx.util import Cm

def add_slide(prs: Presentation, layout_index: int = 1, title: str = "", content: str = "") -> Dict[str, Any]:
    """Adds a new slide to the presentation, automatically filling title and content placeholders if available."""
    if layout_index < 0 or layout_index >= len(prs.slide_layouts):
        # Fallback: layout index 1 is usually "Title and Content"
        layout_index = 1 if len(prs.slide_layouts) > 1 else 0

    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)

    # 1. Fill title if specified
    if title:
        title_set = False
        try:
            if slide.shapes.title:
                slide.shapes.title.text = title
                title_set = True
        except Exception:
            pass

        if not title_set:
            # Search for title placeholder manually
            for shape in slide.shapes:
                if shape.is_placeholder and hasattr(shape, "placeholder_format"):
                    if shape.placeholder_format.type in (1, 3):  # TITLE or CENTER_TITLE
                        shape.text = title
                        title_set = True
                        break

    # 2. Fill content if specified
    if content:
        content_set = False
        # Get reference to title shape to avoid double filling
        title_shape = None
        try:
            title_shape = slide.shapes.title
        except Exception:
            pass

        for shape in slide.shapes:
            if shape == title_shape:
                continue
            if shape.is_placeholder and shape.has_text_frame:
                # Fill placeholder text frame
                tf = shape.text_frame
                tf.clear()
                lines = content.split("\n")
                for i, line in enumerate(lines):
                    if i == 0:
                        tf.paragraphs[0].text = line
                    else:
                        p = tf.add_paragraph()
                        p.text = line
                content_set = True
                break

        # If no placeholder was found/set, add a floating textbox
        if not content_set:
            left = Cm(2.0)
            top = Cm(5.0)
            width = Cm(21.0)
            height = Cm(10.0)
            tb = slide.shapes.add_textbox(left, top, width, height)
            tf = tb.text_frame
            tf.clear()
            lines = content.split("\n")
            for i, line in enumerate(lines):
                if i == 0:
                    tf.paragraphs[0].text = line
                else:
                    p = tf.add_paragraph()
                    p.text = line

    return {
        "slide_index": len(prs.slides) - 1,
        "layout_index": layout_index,
        "layout_name": layout.name,
        "title": title,
        "content_added": bool(content)
    }

def add_text_box(prs: Presentation, slide_index: int, text: str, left_cm: float, top_cm: float, width_cm: float, height_cm: float) -> Dict[str, Any]:
    """Adds a floating textbox at specified coordinates with multiline text formatting."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"Slide index {slide_index} is out of bounds (0 to {len(prs.slides) - 1}).")

    slide = prs.slides[slide_index]
    left = Cm(left_cm)
    top = Cm(top_cm)
    width = Cm(width_cm)
    height = Cm(height_cm)

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()

    lines = text.split("\n")
    for idx, line in enumerate(lines):
        if idx == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line

    # Return shape info
    return {
        "slide_index": slide_index,
        "shape_name": tb.name,
        "shape_id": tb.shape_id,
        "text": text,
        "left_cm": left_cm,
        "top_cm": top_cm,
        "width_cm": width_cm,
        "height_cm": height_cm
    }

def add_image(prs: Presentation, slide_index: int, image_path: str, left_cm: float, top_cm: float, width_cm: Optional[float] = None, height_cm: Optional[float] = None) -> Dict[str, Any]:
    """Adds an image at specified coordinates, optionally preserving native aspect ratio if width/height are omitted."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"Slide index {slide_index} is out of bounds (0 to {len(prs.slides) - 1}).")

    abs_image_path = os.path.abspath(image_path)
    if not os.path.exists(abs_image_path):
        raise FileNotFoundError(f"Image not found at path: {abs_image_path}")

    slide = prs.slides[slide_index]
    left = Cm(left_cm)
    top = Cm(top_cm)
    width = Cm(width_cm) if width_cm is not None else None
    height = Cm(height_cm) if height_cm is not None else None

    picture = slide.shapes.add_picture(abs_image_path, left, top, width=width, height=height)

    return {
        "slide_index": slide_index,
        "shape_name": picture.name,
        "shape_id": picture.shape_id,
        "image_path": abs_image_path,
        "left_cm": left_cm,
        "top_cm": top_cm,
        "width_cm": width_cm,
        "height_cm": height_cm
    }

def add_table(prs: Presentation, slide_index: int, rows: int, cols: int, left_cm: float, top_cm: float, width_cm: float, height_cm: float, data: Optional[List[List[Any]]] = None) -> Dict[str, Any]:
    """Adds a grid table on the slide, optionally pre-populating cell data."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"Slide index {slide_index} is out of bounds (0 to {len(prs.slides) - 1}).")

    slide = prs.slides[slide_index]
    left = Cm(left_cm)
    top = Cm(top_cm)
    width = Cm(width_cm)
    height = Cm(height_cm)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Populate cell texts if data is provided
    if data:
        for r_idx, row_data in enumerate(data):
            if r_idx >= rows:
                break
            for c_idx, cell_value in enumerate(row_data):
                if c_idx >= cols:
                    break
                table.cell(r_idx, c_idx).text = str(cell_value)

    return {
        "slide_index": slide_index,
        "shape_name": table_shape.name,
        "shape_id": table_shape.shape_id,
        "rows": rows,
        "columns": cols,
        "left_cm": left_cm,
        "top_cm": top_cm,
        "width_cm": width_cm,
        "height_cm": height_cm
    }

def list_layouts(prs: Presentation) -> List[Dict[str, Any]]:
    """Returns available slide layouts with their indices and names."""
    layouts = []
    for idx, layout in enumerate(prs.slide_layouts):
        layouts.append({
            "layout_index": idx,
            "layout_name": layout.name
        })
    return layouts
