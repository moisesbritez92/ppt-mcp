import os
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def to_cm(length) -> Optional[float]:
    """Helper to convert EMU length to Centimeters safely."""
    if length is None:
        return None
    try:
        return round(length.cm, 2)
    except Exception:
        # Fallback: 1 cm = 360,000 EMUs
        return round(length / 360000.0, 2)

def get_slide_title(slide) -> str:
    """Finds the title of a slide with robust fallbacks."""
    try:
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title_text = slide.shapes.title.text_frame.text.strip()
            if title_text:
                return title_text
    except Exception:
        pass

    # Fallback 1: Search for placeholders that are title/center-title
    for shape in slide.shapes:
        try:
            if shape.is_placeholder and hasattr(shape, "placeholder_format"):
                # Title placeholder types are typically 1 (TITLE) or 3 (CENTER_TITLE)
                if shape.placeholder_format.type in (1, 3) and shape.has_text_frame:
                    title_text = shape.text_frame.text.strip()
                    if title_text:
                        return title_text
        except Exception:
            pass

    # Fallback 2: Search for any non-empty text shape
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    # Truncate if too long to represent a title
                    return text[:50] + "..." if len(text) > 50 else text
        except Exception:
            pass

    return f"Untitled Slide"

def parse_shape_details(shape) -> Dict[str, Any]:
    """Extracts relevant info from a single shape recursively if group shape."""
    shape_type_name = "UNKNOWN"
    try:
        if hasattr(shape, "shape_type") and shape.shape_type is not None:
            shape_type_name = shape.shape_type.name
    except Exception:
        pass

    details: Dict[str, Any] = {
        "name": shape.name,
        "shape_id": shape.shape_id,
        "type": shape_type_name,
        "left_cm": to_cm(shape.left),
        "top_cm": to_cm(shape.top),
        "width_cm": to_cm(shape.width),
        "height_cm": to_cm(shape.height),
    }

    # Text content
    if shape.has_text_frame:
        details["text"] = shape.text_frame.text

    # Table content
    if shape.has_table:
        table = shape.table
        rows_data = []
        for row_idx, row in enumerate(table.rows):
            row_cells = []
            for col_idx, cell in enumerate(row.cells):
                row_cells.append({
                    "row": row_idx,
                    "col": col_idx,
                    "text": cell.text
                })
            rows_data.append(row_cells)
        details["table"] = {
            "rows_count": len(table.rows),
            "columns_count": len(table.columns),
            "data": rows_data
        }

    # Group shape processing
    if shape_type_name == "GROUP":
        sub_shapes = []
        try:
            for sub_shape in shape.shapes:
                sub_shapes.append(parse_shape_details(sub_shape))
        except Exception:
            pass
        details["sub_shapes"] = sub_shapes

    return details

def get_presentation_structure(prs: Presentation) -> List[Dict[str, Any]]:
    """Lists slides with index, layout_name, title, and shape count."""
    structure = []
    for idx, slide in enumerate(prs.slides):
        layout_name = "Unknown"
        try:
            layout_name = slide.slide_layout.name
        except Exception:
            pass

        structure.append({
            "slide_index": idx,
            "layout_name": layout_name,
            "title": get_slide_title(slide),
            "shapes_count": len(slide.shapes)
        })
    return structure

def read_slide(prs: Presentation, slide_index: int) -> Dict[str, Any]:
    """Gets detailed info about shapes, position, texts, and tables of a slide."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"Slide index {slide_index} is out of bounds (0 to {len(prs.slides) - 1}).")

    slide = prs.slides[slide_index]
    layout_name = "Unknown"
    try:
        layout_name = slide.slide_layout.name
    except Exception:
        pass

    shapes_details = []
    for shape in slide.shapes:
        shapes_details.append(parse_shape_details(shape))

    # Get notes
    notes = ""
    if hasattr(slide, "has_notes_slide") and slide.has_notes_slide:
        try:
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                notes = notes_slide.notes_text_frame.text
        except Exception:
            pass
    elif not hasattr(slide, "has_notes_slide"):
        # Fallback for python-pptx versions without has_notes_slide attribute
        try:
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                notes = notes_slide.notes_text_frame.text
        except Exception:
            pass

    return {
        "slide_index": slide_index,
        "layout_name": layout_name,
        "title": get_slide_title(slide),
        "notes": notes,
        "shapes": shapes_details
    }

def read_slide_notes(prs: Presentation, slide_index: int) -> str:
    """Returns the notes of a slide."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"Slide index {slide_index} is out of bounds (0 to {len(prs.slides) - 1}).")

    slide = prs.slides[slide_index]
    notes = ""
    try:
        if hasattr(slide, "has_notes_slide") and not slide.has_notes_slide:
            return ""
        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            notes = notes_slide.notes_text_frame.text
    except Exception:
        pass
    return notes

def read_full_text(prs: Presentation) -> List[Dict[str, Any]]:
    """Returns all text elements organized by slide."""
    full_text_list = []
    for idx, slide in enumerate(prs.slides):
        texts = []
        
        # Recursive text extractor for shapes
        def extract_text(shapes):
            for shape in shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    texts.append({
                        "shape_name": shape.name,
                        "shape_id": shape.shape_id,
                        "text": shape.text_frame.text
                    })
                if shape.has_table:
                    for row_idx, row in enumerate(shape.table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            if cell.text.strip():
                                texts.append({
                                    "shape_name": f"{shape.name} [Table Row {row_idx}, Col {col_idx}]",
                                    "shape_id": shape.shape_id,
                                    "text": cell.text
                                })
                # Group shape handling
                if hasattr(shape, "shape_type") and shape.shape_type is not None:
                    if shape.shape_type.name == "GROUP":
                        try:
                            extract_text(shape.shapes)
                        except Exception:
                            pass

        extract_text(slide.shapes)
        
        # Include notes text if present
        notes = read_slide_notes(prs, idx)

        full_text_list.append({
            "slide_index": idx,
            "title": get_slide_title(slide),
            "text_elements": texts,
            "notes": notes
        })
    return full_text_list

def list_slides(prs: Presentation) -> List[Dict[str, Any]]:
    """Returns a quick summary list of all slides."""
    return get_presentation_structure(prs)

def search_text(prs: Presentation, query: str, case_sensitive: bool = False) -> List[Dict[str, Any]]:
    """Searches for a text query in all slides, returning matches with context."""
    results = []
    query_str = query if case_sensitive else query.lower()

    for idx, slide in enumerate(prs.slides):
        slide_title = get_slide_title(slide)

        def search_shapes(shapes):
            for shape in shapes:
                # Text Frame Search
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    text_cmp = text if case_sensitive else text.lower()
                    if query_str in text_cmp:
                        # Find line/snippet
                        snippet = ""
                        lines = text.split("\n")
                        for line in lines:
                            line_cmp = line if case_sensitive else line.lower()
                            if query_str in line_cmp:
                                snippet = line.strip()
                                break
                        if not snippet and lines:
                            snippet = lines[0].strip()

                        results.append({
                            "slide_index": idx,
                            "slide_title": slide_title,
                            "shape_id": shape.shape_id,
                            "shape_name": shape.name,
                            "type": "text_frame",
                            "snippet": snippet or text[:100]
                        })

                # Table Search
                if shape.has_table:
                    for r_idx, row in enumerate(shape.table.rows):
                        for c_idx, cell in enumerate(row.cells):
                            cell_text = cell.text
                            cell_cmp = cell_text if case_sensitive else cell_text.lower()
                            if query_str in cell_cmp:
                                results.append({
                                    "slide_index": idx,
                                    "slide_title": slide_title,
                                    "shape_id": shape.shape_id,
                                    "shape_name": f"{shape.name} (Row {r_idx}, Col {c_idx})",
                                    "type": "table_cell",
                                    "snippet": cell_text.strip()
                                })

                # Notes Search
                # We handle notes separately below, but this is for shape search.
                if hasattr(shape, "shape_type") and shape.shape_type is not None:
                    if shape.shape_type.name == "GROUP":
                        try:
                            search_shapes(shape.shapes)
                        except Exception:
                            pass

        search_shapes(slide.shapes)

        # Search Notes
        notes = read_slide_notes(prs, idx)
        notes_cmp = notes if case_sensitive else notes.lower()
        if query_str in notes_cmp:
            # Find snippet
            snippet = ""
            lines = notes.split("\n")
            for line in lines:
                line_cmp = line if case_sensitive else line.lower()
                if query_str in line_cmp:
                    snippet = line.strip()
                    break
            if not snippet and lines:
                snippet = lines[0].strip()

            results.append({
                "slide_index": idx,
                "slide_title": slide_title,
                "shape_id": None,
                "shape_name": "Notes",
                "type": "notes",
                "snippet": snippet or notes[:100]
            })

    return results
