import os
import re
import io
import copy
from typing import Any, Dict, List, Union, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def find_shape(slide, name_or_index) -> Any:
    """Finds a shape on a slide by name, ID, or index."""
    # Try index first if it looks like an integer
    try:
        idx = int(name_or_index)
        if 0 <= idx < len(slide.shapes):
            return slide.shapes[idx]
    except (ValueError, TypeError):
        pass

    # Search by name or shape_id
    for shape in slide.shapes:
        if shape.name == name_or_index:
            return shape
        if str(shape.shape_id) == str(name_or_index):
            return shape
            
    raise KeyError(
        f"Shape '{name_or_index}' not found in slide. "
        f"Available shapes: {[s.name for s in slide.shapes]}"
    )

def edit_shape_text(prs: Presentation, slide_index: int, shape_name_or_index: Union[str, int], new_text: str) -> Dict[str, Any]:
    """Edits the text of a shape, splitting newlines into paragraphs to preserve structure."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"Slide index {slide_index} is out of bounds (0 to {len(prs.slides) - 1}).")

    slide = prs.slides[slide_index]
    shape = find_shape(slide, shape_name_or_index)

    if not shape.has_text_frame:
        raise ValueError(f"Shape '{shape.name}' (type: {shape.shape_type}) does not support text frames.")

    tf = shape.text_frame
    tf.clear()  # clears all paragraphs except the first empty one
    
    lines = new_text.split("\n")
    for idx, line in enumerate(lines):
        if idx == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line

    return {
        "slide_index": slide_index,
        "shape_name": shape.name,
        "shape_id": shape.shape_id,
        "new_text": new_text
    }

def edit_slide_title(prs: Presentation, slide_index: int, new_title: str) -> Dict[str, Any]:
    """Edits the title of a slide, using shapes.title or a title placeholder fallback."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"Slide index {slide_index} is out of bounds (0 to {len(prs.slides) - 1}).")

    slide = prs.slides[slide_index]

    # Try setting slide.shapes.title
    title_shape = None
    try:
        title_shape = slide.shapes.title
    except Exception:
        pass

    if title_shape:
        title_shape.text = new_title
        return {
            "slide_index": slide_index,
            "title_source": "shapes.title",
            "new_title": new_title
        }

    # Fallback: Find a title placeholder manually
    for shape in slide.shapes:
        if shape.is_placeholder and hasattr(shape, "placeholder_format"):
            # Title placeholders are 1 (TITLE) or 3 (CENTER_TITLE)
            if shape.placeholder_format.type in (1, 3):
                shape.text = new_title
                return {
                    "slide_index": slide_index,
                    "title_source": f"placeholder (index: {shape.placeholder_format.idx})",
                    "new_title": new_title
                }

    # Second fallback: Add a title textbox if none exists
    # Or just raise error. Let's raise ValueError to let the user know they can create one.
    raise ValueError(
        f"No title shape or title placeholder found on slide {slide_index}. "
        "Please use 'edit_shape_text' or add a text box."
    )

def edit_slide_notes(prs: Presentation, slide_index: int, new_text: str) -> Dict[str, Any]:
    """Edits/sets the presenter notes of a slide."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"Slide index {slide_index} is out of bounds (0 to {len(prs.slides) - 1}).")

    slide = prs.slides[slide_index]
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = new_text

    return {
        "slide_index": slide_index,
        "notes": new_text
    }

def _replace_in_text_frame(tf, search: str, replace: str, case_sensitive: bool) -> bool:
    """Replaces text in a text frame. Tries run-level first, falls back to paragraph-level."""
    replaced_any = False
    search_term = search if case_sensitive else search.lower()

    for paragraph in tf.paragraphs:
        # 1. Run-level search & replace (preserves styles if word is fully inside a run)
        run_replaced = False
        for run in paragraph.runs:
            run_text_cmp = run.text if case_sensitive else run.text.lower()
            if search_term in run_text_cmp:
                if case_sensitive:
                    run.text = run.text.replace(search, replace)
                else:
                    pattern = re.compile(re.escape(search), re.IGNORECASE)
                    run.text = pattern.sub(replace, run.text)
                run_replaced = True
                replaced_any = True

        # 2. Paragraph-level fallback (only if search phrase is split across multiple runs)
        if not run_replaced:
            para_text_cmp = paragraph.text if case_sensitive else paragraph.text.lower()
            if search_term in para_text_cmp:
                # Replaced at paragraph level, which merges runs into a single run
                if case_sensitive:
                    paragraph.text = paragraph.text.replace(search, replace)
                else:
                    pattern = re.compile(re.escape(search), re.IGNORECASE)
                    paragraph.text = pattern.sub(replace, paragraph.text)
                replaced_any = True

    return replaced_any

def replace_text(prs: Presentation, search: str, replace: str, case_sensitive: bool = False) -> Dict[str, Any]:
    """Performs global search and replace across text frames, tables, and notes."""
    replace_count = 0
    slides_affected = set()

    for idx, slide in enumerate(prs.slides):
        slide_replaced = False

        def walk_and_replace(shapes):
            nonlocal replace_count, slide_replaced
            for shape in shapes:
                # Textbox
                if shape.has_text_frame:
                    if _replace_in_text_frame(shape.text_frame, search, replace, case_sensitive):
                        replace_count += 1
                        slide_replaced = True

                # Table
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text:
                                if _replace_in_text_frame(cell.text_frame, search, replace, case_sensitive):
                                    replace_count += 1
                                    slide_replaced = True

                # Group shape recursion
                if hasattr(shape, "shape_type") and shape.shape_type is not None:
                    if shape.shape_type.name == "GROUP":
                        try:
                            walk_and_replace(shape.shapes)
                        except Exception:
                            pass

        walk_and_replace(slide.shapes)

        # Slide Notes
        if hasattr(slide, "has_notes_slide") and slide.has_notes_slide:
            try:
                if _replace_in_text_frame(slide.notes_slide.notes_text_frame, search, replace, case_sensitive):
                    replace_count += 1
                    slide_replaced = True
            except Exception:
                pass
        elif not hasattr(slide, "has_notes_slide"):
            try:
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame:
                    if _replace_in_text_frame(notes_slide.notes_text_frame, search, replace, case_sensitive):
                        replace_count += 1
                        slide_replaced = True
            except Exception:
                pass

        if slide_replaced:
            slides_affected.add(idx)

    return {
        "search_term": search,
        "replace_term": replace,
        "occurrences_replaced": replace_count,
        "slides_affected": sorted(list(slides_affected))
    }

def delete_slide(prs: Presentation, slide_index: int) -> Dict[str, Any]:
    """Deletes a slide from the presentation using python-pptx XML manipulation."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"Slide index {slide_index} is out of bounds (0 to {len(prs.slides) - 1}).")

    slide = prs.slides[slide_index]
    
    # 1. Remove slide from list of slide IDs in presentation XML
    slide_id = prs.slides._sldIdLst[slide_index]
    prs.slides._sldIdLst.remove(slide_id)
    
    # 2. Drop slide relationships
    for r_id, rel in prs.part.rels.items():
        if rel.target_part == slide.part:
            prs.part.drop_rel(r_id)
            break

    return {
        "deleted_slide_index": slide_index,
        "remaining_slides_count": len(prs.slides)
    }

def duplicate_slide(prs: Presentation, slide_index: int) -> Dict[str, Any]:
    """Duplicates a slide, copying placeholders, text shapes, tables, and images."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"Slide index {slide_index} is out of bounds (0 to {len(prs.slides) - 1}).")

    source_slide = prs.slides[slide_index]
    
    # 1. Add new slide with same layout at the end
    dest_slide = prs.slides.add_slide(source_slide.slide_layout)
    
    # 2. Match placeholders and copy their texts
    source_placeholders = {s.placeholder_format.idx: s for s in source_slide.shapes if s.is_placeholder}
    dest_placeholders = {s.placeholder_format.idx: s for s in dest_slide.shapes if s.is_placeholder}
    
    for idx_key, s_shape in source_placeholders.items():
        if idx_key in dest_placeholders:
            d_shape = dest_placeholders[idx_key]
            if s_shape.has_text_frame and s_shape.text_frame.text.strip():
                d_shape.text = s_shape.text_frame.text

    # 3. Recreate non-placeholder shapes
    for shape in source_slide.shapes:
        if shape.is_placeholder:
            continue
        
        # Textbox or general AutoShape
        if shape.has_text_frame and not shape.has_table:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                new_shape = dest_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                new_shape.text_frame.text = shape.text_frame.text
            else:
                try:
                    new_shape = dest_slide.shapes.add_shape(
                        shape.shape_type, shape.left, shape.top, shape.width, shape.height
                    )
                    if new_shape.has_text_frame:
                        new_shape.text_frame.text = shape.text_frame.text
                except Exception:
                    # Fallback to textbox
                    new_shape = dest_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                    new_shape.text_frame.text = shape.text_frame.text

        # Table
        elif shape.has_table:
            s_table = shape.table
            rows = len(s_table.rows)
            cols = len(s_table.columns)
            new_shape = dest_slide.shapes.add_table(
                rows, cols, shape.left, shape.top, shape.width, shape.height
            )
            d_table = new_shape.table
            for r in range(rows):
                for c in range(cols):
                    d_table.cell(r, c).text = s_table.cell(r, c).text

        # Picture
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image_blob = shape.image.blob
                img_stream = io.BytesIO(image_blob)
                dest_slide.shapes.add_picture(
                    img_stream, shape.left, shape.top, shape.width, shape.height
                )
            except Exception:
                pass

    # 4. Copy slide notes
    try:
        if hasattr(source_slide, "has_notes_slide") and source_slide.has_notes_slide:
            dest_slide.notes_slide.notes_text_frame.text = source_slide.notes_slide.notes_text_frame.text
        elif not hasattr(source_slide, "has_notes_slide"):
            if source_slide.notes_slide and source_slide.notes_slide.notes_text_frame:
                dest_slide.notes_slide.notes_text_frame.text = source_slide.notes_slide.notes_text_frame.text
    except Exception:
        pass

    return {
        "duplicated_slide_index": slide_index,
        "new_slide_index": len(prs.slides) - 1,
        "total_slides": len(prs.slides)
    }

def move_slide(prs: Presentation, from_index: int, to_index: int) -> Dict[str, Any]:
    """Moves a slide from from_index to to_index by reordering sldIdLst."""
    if from_index < 0 or from_index >= len(prs.slides):
        raise IndexError(f"from_index {from_index} is out of bounds (0 to {len(prs.slides) - 1}).")
    if to_index < 0 or to_index >= len(prs.slides):
        raise IndexError(f"to_index {to_index} is out of bounds (0 to {len(prs.slides) - 1}).")

    sldIdLst = prs.slides._sldIdLst
    slide_id_element = sldIdLst[from_index]
    
    sldIdLst.remove(slide_id_element)
    sldIdLst.insert(to_index, slide_id_element)

    return {
        "from_index": from_index,
        "to_index": to_index,
        "total_slides": len(prs.slides)
    }
