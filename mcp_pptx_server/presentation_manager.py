import os
import json
from dataclasses import dataclass
from typing import Dict, Optional, List, Any
from pptx import Presentation

@dataclass
class PresentationEntry:
    presentation: Presentation
    file_path: Optional[str] = None

class PresentationManager:
    def __init__(self):
        self._presentations: Dict[str, PresentationEntry] = {}

    def create(self, prs_id: str, title: str = "") -> PresentationEntry:
        """Creates a new empty presentation in memory with a title slide if title is provided."""
        prs = Presentation()
        if title:
            # Layout 0 is typically the title slide in default templates
            if len(prs.slide_layouts) > 0:
                title_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_layout)
                # Find and set title
                if slide.shapes.title:
                    slide.shapes.title.text = title
                else:
                    # Fallback if no explicit title shape exists, try setting on any placeholder
                    for shape in slide.shapes:
                        if shape.is_placeholder:
                            shape.text = title
                            break
        
        entry = PresentationEntry(presentation=prs, file_path=None)
        self._presentations[prs_id] = entry
        return entry

    def open(self, prs_id: str, file_path: str) -> PresentationEntry:
        """Opens an existing presentation from disk."""
        abs_path = os.path.abspath(file_path)
        if not os.path.exists(abs_path):
            raise FileNotFoundError(f"File not found: {abs_path}")
        
        prs = Presentation(abs_path)
        entry = PresentationEntry(presentation=prs, file_path=abs_path)
        self._presentations[prs_id] = entry
        return entry

    def save(self, prs_id: str, file_path: Optional[str] = None) -> str:
        """Saves an open presentation. If file_path is not specified, saves to its original path."""
        entry = self.get(prs_id)
        
        save_path = file_path or entry.file_path
        if not save_path:
            raise ValueError(f"No file path specified for presentation '{prs_id}' and it was not loaded from a file.")
        
        abs_path = os.path.abspath(save_path)
        # Ensure parent directory exists
        os.makedirs(os.path.dirname(abs_path), exist_ok=True)
        
        entry.presentation.save(abs_path)
        entry.file_path = abs_path
        return abs_path

    def close(self, prs_id: str) -> None:
        """Closes a presentation, freeing it from memory."""
        if prs_id in self._presentations:
            del self._presentations[prs_id]
        else:
            raise KeyError(
                f"Presentation '{prs_id}' not found. "
                f"Available presentations: {list(self._presentations.keys())}"
            )

    def list_presentations(self) -> List[Dict[str, Any]]:
        """Lists all open presentations in memory."""
        results = []
        for prs_id, entry in self._presentations.items():
            results.append({
                "prs_id": prs_id,
                "file_path": entry.file_path or "In-Memory (unsaved)",
                "slides_count": len(entry.presentation.slides)
            })
        return results

    def get(self, prs_id: str) -> PresentationEntry:
        """Gets a presentation entry by prs_id or raises a descriptive error."""
        if prs_id not in self._presentations:
            raise KeyError(
                f"Presentation '{prs_id}' not found. "
                f"Available presentations: {list(self._presentations.keys())}. "
                "Please call 'create_presentation' or 'open_presentation' first."
            )
        return self._presentations[prs_id]
